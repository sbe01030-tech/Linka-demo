"""
Linka — 1페이지 앱 소개 PPT 생성.
슬로건: Worth Work, Work Worth
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── Colors (Linka 디자인 시스템) ─────────────────────────
ACCENT      = RGBColor(0x00, 0xC8, 0x53)   # 브랜드 그린
DARK        = RGBColor(0x11, 0x18, 0x27)
GRAY        = RGBColor(0x6B, 0x72, 0x80)
GRAY_LIGHT  = RGBColor(0x9C, 0xA3, 0xAF)
HELPER      = RGBColor(0xF5, 0x9E, 0x0B)   # amber
DRIVER      = RGBColor(0x63, 0x66, 0xF1)   # indigo
ERRAND      = RGBColor(0xF9, 0x73, 0x16)   # orange
SECTION     = RGBColor(0xF9, 0xFA, 0xFB)
BORDER      = RGBColor(0xE5, 0xE7, 0xEB)

# ── Slide setup (16:9 widescreen) ─────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# ── Helpers ─────────────────────────────────────────
def text_box(slide, x, y, w, h, text, *,
             size=14, bold=False, color=DARK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font="맑은 고딕"):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb

def rect(slide, x, y, w, h, fill, *, line=None, radius=False):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    return shape

# ── 1. 헤더 (로고 + 위치 태그) ─────────────────────────
text_box(slide, 0.6, 0.4, 3.0, 0.6,
         "Linka", size=36, bold=True, color=DARK, font="Nunito")

# 우상단 location tag
loc_box = rect(slide, 11.2, 0.45, 1.7, 0.4, SECTION, line=BORDER, radius=True)
text_box(slide, 11.2, 0.45, 1.7, 0.4,
         "🇮🇩 Jakarta · 2026", size=10, color=GRAY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ── 2. 슬로건 (메인) ─────────────────────────────────
text_box(slide, 0.6, 1.3, 12.1, 1.2,
         "Worth Work, Work Worth",
         size=54, bold=True, color=ACCENT, align=PP_ALIGN.LEFT, font="Nunito")

# 한 줄 카피
text_box(slide, 0.6, 2.5, 12.1, 0.5,
         "인도네시아 가정의 손길을 잇는 신뢰 기반 매칭 플랫폼",
         size=16, color=GRAY)

# 구분선
rect(slide, 0.6, 3.15, 12.1, 0.02, BORDER)

# ── 3. 3대 서비스 (가로 3등분) ───────────────────────
SERVICES = [
    ("가사도우미",  "ART · 청소 · 요리 · 육아",      HELPER,  "🏠"),
    ("드라이버",    "대리 · 일일 · 시간제 · 공항",     DRIVER,  "🚗"),
    ("심부름",      "누구나 등록 · 누구나 지원",       ERRAND,  "📦"),
]
card_w  = 4.0
card_h  = 1.6
gap     = 0.15
total_w = card_w * 3 + gap * 2
start_x = (13.333 - total_w) / 2
y0      = 3.4

for i, (name, sub, color, emoji) in enumerate(SERVICES):
    x = start_x + i * (card_w + gap)
    # 카드 배경
    rect(slide, x, y0, card_w, card_h, SECTION, line=BORDER, radius=True)
    # 좌측 컬러 바
    rect(slide, x, y0, 0.12, card_h, color)
    # 이모지
    text_box(slide, x + 0.3, y0 + 0.15, 0.8, 0.6, emoji,
             size=30, anchor=MSO_ANCHOR.MIDDLE)
    # 이름
    text_box(slide, x + 1.1, y0 + 0.2, card_w - 1.3, 0.5, name,
             size=20, bold=True, color=DARK)
    # 서브
    text_box(slide, x + 1.1, y0 + 0.75, card_w - 1.3, 0.7, sub,
             size=12, color=GRAY)

# ── 4. 차별점 3개 ────────────────────────────────────
DIFFS = [
    ("KYC + Linka 온도",
     "인니식 본인인증(NIK·KTP) + 평점 통합 → 신뢰 가시화"),
    ("통합 슈퍼앱 전략",
     "가사·드라이버·심부름을 한 앱에 → 가족 단위 LTV ↑"),
    ("프리미엄·주재원 교두보",
     "자카르타 중상층 + 한국·일본 expat 우선 타겟"),
]
y1 = 5.3
text_box(slide, 0.6, y1, 12.1, 0.4, "우리의 차별점",
         size=14, bold=True, color=ACCENT)

for i, (title, desc) in enumerate(DIFFS):
    col_w = 4.0
    x = 0.6 + i * (col_w + 0.15)
    # number badge
    rect(slide, x, y1 + 0.55, 0.35, 0.35, ACCENT, radius=True)
    text_box(slide, x, y1 + 0.55, 0.35, 0.35, str(i + 1),
             size=14, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # title
    text_box(slide, x + 0.5, y1 + 0.5, col_w - 0.5, 0.4, title,
             size=13, bold=True, color=DARK)
    # desc
    text_box(slide, x + 0.5, y1 + 0.9, col_w - 0.5, 0.7, desc,
             size=11, color=GRAY)

# ── 5. 푸터 ──────────────────────────────────────────
text_box(slide, 0.6, 7.0, 8.0, 0.3,
         "2026 시연 단계 · Pre-launch · React Native + Java/AWS Backend",
         size=10, color=GRAY_LIGHT)

text_box(slide, 9.0, 7.0, 4.0, 0.3,
         "Park · max@bigbandent.com",
         size=10, color=GRAY_LIGHT, align=PP_ALIGN.RIGHT)

# ── Save ────────────────────────────────────────────
out = "/Users/bigband/Linka-app/Linka-app/docs/Linka-onepager.pptx"
prs.save(out)
print(f"OK — {out}")
