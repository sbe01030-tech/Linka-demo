"""
Linka 백엔드 기획안 PPT v2 생성기.
- 폰 비율(9:19.5) 세로 스크린샷 placeholder
- 백엔드 연결 항목 확장 (즐겨찾기·신고·추천·실시간·이벤트·지오코딩 등)
실행: /tmp/linka-ppt/bin/python docs/backend-spec/generate_ppt_v2.py
출력: docs/backend-spec/Linka-backend-spec-v2.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Constants ────────────────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

KO_FONT = "Apple SD Gothic Neo"
EN_FONT = "SF Pro Text"
MONO_FONT = "SF Mono"

BRAND       = RGBColor(0x00, 0xC8, 0x53)
INK         = RGBColor(0x10, 0x14, 0x18)
SUB         = RGBColor(0x55, 0x60, 0x6C)
LINE        = RGBColor(0xE3, 0xE6, 0xEA)
BG_SOFT     = RGBColor(0xF5, 0xF7, 0xF8)
HELPER_CLR  = RGBColor(0xF5, 0x9E, 0x0B)
DRIVER_CLR  = RGBColor(0x63, 0x66, 0xF1)
WARN        = RGBColor(0xEF, 0x44, 0x44)
INFO        = RGBColor(0x3B, 0x82, 0xF6)
PURPLE      = RGBColor(0x8B, 0x5C, 0xF6)


def set_font(run, *, size=14, bold=False, color=INK, font=KO_FONT):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_text(slide, text, *, left, top, width, height,
             size=14, bold=False, color=INK, align=PP_ALIGN.LEFT, font=KO_FONT):
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_font(run, size=size, bold=bold, color=color, font=font)
    return tx


def add_bullets(slide, items, *, left, top, width, height,
                size=11, color=INK, line_spacing=1.35, font=KO_FONT):
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = "•  " + item
        set_font(run, size=size, color=color, font=font)
    return tx


def add_rect(slide, left, top, width, height, *, fill=BG_SOFT, line=None, radius=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.shadow.inherit = False
    if radius:
        shape.adjustments[0] = radius
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.5)
    return shape


def add_pill(slide, text, *, left, top, fill, color=RGBColor(0xFF, 0xFF, 0xFF), size=10, width=None):
    width = width or Inches(1.4)
    height = Inches(0.32)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.shadow.inherit = False
    shape.adjustments[0] = 0.5
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    set_font(run, size=size, bold=True, color=color)
    return shape


def add_header(slide, eyebrow, title, *, page_label=""):
    add_text(slide, eyebrow, left=Inches(0.6), top=Inches(0.4),
             width=Inches(9), height=Inches(0.3), size=11, bold=True, color=BRAND)
    add_text(slide, title, left=Inches(0.6), top=Inches(0.72),
             width=Inches(11), height=Inches(0.6), size=24, bold=True, color=INK)
    add_rect(slide, Inches(0.6), Inches(1.40), Inches(12.13), Emu(9525), fill=LINE)
    if page_label:
        add_text(slide, page_label, left=Inches(11.0), top=Inches(0.45),
                 width=Inches(1.7), height=Inches(0.3), size=10, color=SUB,
                 align=PP_ALIGN.RIGHT)


def add_footer(slide, page_no, total):
    add_text(slide, "Linka  ·  Backend Spec  ·  v0.2  ·  Map screen",
             left=Inches(0.6), top=Inches(7.10),
             width=Inches(7), height=Inches(0.3), size=9, color=SUB)
    add_text(slide, f"{page_no} / {total}",
             left=Inches(11.0), top=Inches(7.10),
             width=Inches(1.7), height=Inches(0.3), size=9, color=SUB,
             align=PP_ALIGN.RIGHT)


# ── 폰 비율 placeholder (9:19.5) ───────────────────────────
PHONE_W = Inches(2.55)
PHONE_H = Inches(5.52)


def add_phone_placeholder(slide, left, top, *, label, sub=None, w=PHONE_W, h=PHONE_H):
    # 외부 프레임 (폰 케이스)
    add_rect(slide, left, top, w, h, fill=RGBColor(0x10, 0x14, 0x18),
             line=None, radius=0.10)
    # 내부 화면
    inner_pad = Inches(0.07)
    add_rect(slide, left + inner_pad, top + inner_pad,
             w - inner_pad * 2, h - inner_pad * 2,
             fill=BG_SOFT, line=None, radius=0.06)
    # 노치 (얇은 핀)
    notch_w = Inches(0.7)
    notch_h = Inches(0.13)
    add_rect(slide,
             left + (w - notch_w) / 2, top + Inches(0.18),
             notch_w, notch_h,
             fill=RGBColor(0x10, 0x14, 0x18), line=None, radius=0.5)
    # 안내 텍스트
    add_text(slide, "📷  스크린샷",
             left=left, top=top + h / 2 - Inches(0.4),
             width=w, height=Inches(0.4), size=11, color=SUB,
             align=PP_ALIGN.CENTER)
    if label:
        add_text(slide, label,
                 left=left, top=top + h / 2 + Inches(0.05),
                 width=w, height=Inches(0.4), size=10, bold=True, color=INK,
                 align=PP_ALIGN.CENTER)
    if sub:
        add_text(slide, sub,
                 left=left, top=top + h / 2 + Inches(0.40),
                 width=w, height=Inches(0.4), size=8, color=SUB,
                 align=PP_ALIGN.CENTER)


def add_phone_caption(slide, left, top, *, title, body):
    """폰 placeholder 아래의 캡션 + 백엔드 연결 메모."""
    add_text(slide, title,
             left=left, top=top, width=PHONE_W, height=Inches(0.3),
             size=10, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_text(slide, body,
             left=left - Inches(0.05), top=top + Inches(0.30),
             width=PHONE_W + Inches(0.1), height=Inches(0.9),
             size=8, color=SUB, align=PP_ALIGN.CENTER)


# ── API 카드 ──────────────────────────────────────────────
def add_api_card(slide, *, left, top, width, height,
                 method, path, title, body_lines,
                 method_color=BRAND, hint=None):
    add_rect(slide, left, top, width, height, fill=BG_SOFT, line=LINE, radius=0.05)
    add_pill(slide, method,
             left=left + Inches(0.18), top=top + Inches(0.18),
             fill=method_color, size=9, width=Inches(0.85))
    add_text(slide, path,
             left=left + Inches(1.10), top=top + Inches(0.18),
             width=width - Inches(1.3), height=Inches(0.32),
             size=11, bold=True, color=INK, font=MONO_FONT)
    add_text(slide, title,
             left=left + Inches(0.18), top=top + Inches(0.55),
             width=width - Inches(0.36), height=Inches(0.3),
             size=10, color=SUB)
    add_bullets(slide, body_lines,
                left=left + Inches(0.18), top=top + Inches(0.85),
                width=width - Inches(0.36), height=height - Inches(1.05),
                size=9, color=INK, line_spacing=1.30)
    if hint:
        add_text(slide, hint,
                 left=left + Inches(0.18), top=top + height - Inches(0.32),
                 width=width - Inches(0.36), height=Inches(0.25),
                 size=8, color=WARN)


# ── Build presentation ──────────────────────────────────────
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]


def new_slide():
    s = prs.slides.add_slide(blank_layout)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=RGBColor(0xFF, 0xFF, 0xFF))
    return s


# ── Slides ─────────────────────────────────────────────────
def s01_cover():
    s = new_slide()
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=RGBColor(0xFA, 0xFB, 0xFC))
    add_rect(s, 0, 0, Inches(0.35), SLIDE_H, fill=BRAND)
    add_text(s, "BACKEND SPEC  ·  v0.2", left=Inches(0.9), top=Inches(1.5),
             width=Inches(8), height=Inches(0.4), size=12, bold=True, color=BRAND)
    add_text(s, "Linka 백엔드 기획안", left=Inches(0.9), top=Inches(1.95),
             width=Inches(11), height=Inches(1.0), size=42, bold=True, color=INK)
    add_text(s,
             "React Native(Expo) 기반 인도네시아 타겟 가사·운전·심부름 중개 앱",
             left=Inches(0.9), top=Inches(3.05),
             width=Inches(11), height=Inches(0.5), size=15, color=SUB)
    # 변경점
    add_rect(s, Inches(0.9), Inches(4.0), Inches(11.5), Inches(2.4),
             fill=BG_SOFT, line=LINE, radius=0.04)
    add_text(s, "v0.2 변경점", left=Inches(1.15), top=Inches(4.20),
             width=Inches(11), height=Inches(0.3), size=12, bold=True, color=BRAND)
    add_bullets(s, [
        "스크린샷 자리를 폰 비율(9:19.5) 세로 placeholder로 교체",
        "백엔드 연결 항목 확장:  즐겨찾기 · 신고 · 차단 · 추천 · 최근 본 · 검색 이력 · 가용 시간",
        "지오코딩 · 활성 지역 · 실시간 가용 상태(WS/SSE) · 분석 이벤트 로깅 · 요금 정책 추가",
        "API 카드 디자인 통일, 메서드 색상 구분 (GET/POST/DELETE/WS)",
        "데이터 모델: Partner 외 보조 엔티티(Activity/Area/Favorite/Report) 추가",
    ], left=Inches(1.15), top=Inches(4.55),
       width=Inches(11.0), height=Inches(1.7), size=11, color=INK, line_spacing=1.4)
    add_text(s, "v0.2  ·  2026-04-28", left=Inches(0.9), top=Inches(6.7),
             width=Inches(6), height=Inches(0.3), size=11, color=SUB)


def s02_index():
    s = new_slide()
    add_header(s, "01  ·  INDEX", "전체 화면 목차")
    cols = [
        ("인증 (Auth)", [
            "01  스플래시", "02  환영 화면 (역할 선택)",
            "03  로그인", "04  회원가입", "05  약관",
        ]),
        ("고객 (Customer)", [
            "10  홈", "11  지도 (파트너 탐색) ← 우선",
            "12  헬퍼 검색", "13  헬퍼 상세",
            "14  예약", "15  주문 내역",
            "16  리뷰 작성", "17  프로필", "18  알림",
        ]),
        ("헬퍼 / 드라이버 / 심부름", [
            "20  헬퍼 홈", "21  헬퍼 주문", "22  헬퍼 프로필",
            "30  드라이버 게시판", "31  드라이버 상세",
            "40  심부름 게시판", "41  심부름 등록", "42  심부름 상세",
        ]),
        ("채팅 / 공통", [
            "50  채팅 목록", "51  채팅 상세",
            "60  프로필 수정", "61  도움말 / FAQ",
            "",
            "※ 커뮤니티 탭은 본 기획안에서 제외",
        ]),
    ]
    col_w = Inches(2.95)
    gap   = Inches(0.15)
    left0 = Inches(0.6)
    top0  = Inches(1.70)
    for i, (title, items) in enumerate(cols):
        left = left0 + (col_w + gap) * i
        add_text(s, title, left=left, top=top0,
                 width=col_w, height=Inches(0.4), size=13, bold=True, color=BRAND)
        add_bullets(s, items, left=left, top=top0 + Inches(0.45),
                    width=col_w, height=Inches(4.5), size=11, color=INK)
    add_footer(s, 2, TOTAL)


def s03_common():
    s = new_slide()
    add_header(s, "02  ·  COMMON", "공통 사항")
    cards = [
        ("사용자 역할", [
            "customer  · 서비스 요청자",
            "helper      · 가사도우미 (ART)",
            "driver       · 운전자 (대리/일일)",
            "errand      · 심부름꾼 (예정)",
            "→ 멀티 역할 허용 여부 정책 결정",
        ]),
        ("다국어", [
            "지원: id (기본) · en · ko · zh · ja",
            "사용자 입력은 원문 저장",
            "카테고리·시스템 메시지만 번역 키 보유",
        ]),
        ("통화 / 지역", [
            "통화: IDR (Rp) · 정수 단위 저장",
            "기본: 자카르타 (Kemang / Kebayoran Baru)",
            "좌표: WGS84 (lat/lng)",
        ]),
        ("인증 / 미디어", [
            "이메일 또는 전화번호 (정책 필요)",
            "JWT 기반 세션 + Refresh",
            "약관 동의 이력 보관",
            "사진·첨부: Pre-signed URL 권장",
        ]),
    ]
    box_w = Inches(5.95)
    box_h = Inches(2.55)
    positions = [
        (Inches(0.6),  Inches(1.70)),
        (Inches(6.78), Inches(1.70)),
        (Inches(0.6),  Inches(4.40)),
        (Inches(6.78), Inches(4.40)),
    ]
    for (title, items), (left, top) in zip(cards, positions):
        add_rect(s, left, top, box_w, box_h, fill=BG_SOFT, line=LINE, radius=0.04)
        add_text(s, title, left=left + Inches(0.3), top=top + Inches(0.25),
                 width=box_w - Inches(0.6), height=Inches(0.4),
                 size=14, bold=True, color=INK)
        add_bullets(s, items, left=left + Inches(0.3), top=top + Inches(0.7),
                    width=box_w - Inches(0.6), height=box_h - Inches(0.85),
                    size=11, color=SUB)
    add_footer(s, 3, TOTAL)


def s04_map_overview():
    s = new_slide()
    add_header(s, "11  ·  CUSTOMER  ·  MAP", "지도 화면 — 개요", page_label="지도 1/12")
    add_text(s, "목적", left=Inches(0.6), top=Inches(1.70),
             width=Inches(6), height=Inches(0.4), size=14, bold=True, color=BRAND)
    add_text(s,
             "지도 위에서 주변 헬퍼·드라이버를 시각적으로 탐색하고, 필터·정렬을 통해 원하는 파트너를 찾아 상세 진입한다.",
             left=Inches(0.6), top=Inches(2.10),
             width=Inches(6), height=Inches(1.2), size=12, color=INK)

    add_text(s, "진입 경로", left=Inches(0.6), top=Inches(3.30),
             width=Inches(6), height=Inches(0.4), size=14, bold=True, color=BRAND)
    add_bullets(s, [
        "하단 탭바 → \"지도(Map)\"",
        "홈 화면 → 서비스 카드 / \"지도에서 보기\"",
        "라우트 파라미터:  expanded?: boolean",
        "라우트 파라미터:  serviceType?: regular | onetime | live-in",
    ], left=Inches(0.6), top=Inches(3.70),
       width=Inches(6), height=Inches(2.2), size=11, color=INK)

    add_text(s, "핵심 지표 (백엔드 모니터링 권장)",
             left=Inches(0.6), top=Inches(5.85),
             width=Inches(6), height=Inches(0.4), size=12, bold=True, color=PURPLE)
    add_bullets(s, [
        "검색 응답 시간 P95 < 400ms",
        "마커→상세 전환율 / 상세→예약 전환율",
        "검색 결과 0건 비율 (필터 UX 개선 신호)",
    ], left=Inches(0.6), top=Inches(6.20),
       width=Inches(6), height=Inches(0.85), size=10, color=SUB, line_spacing=1.3)

    # 우: 파트너 두 종류
    add_rect(s, Inches(7.0), Inches(1.70), Inches(5.7), Inches(5.3),
             fill=BG_SOFT, line=LINE, radius=0.04)
    add_text(s, "두 종류의 파트너", left=Inches(7.3), top=Inches(1.90),
             width=Inches(5.1), height=Inches(0.4), size=14, bold=True, color=INK)
    add_pill(s, "HELPER  ·  헬퍼", left=Inches(7.3), top=Inches(2.45), fill=HELPER_CLR)
    add_bullets(s, [
        "가사도우미 (ART)",
        "스킬: Beberes / Masak / Cuci / Setrika / Deep Cleaning / Masak Sehat",
        "정기 / 단기 / 상주",
        "검증·경력·평점 기반 매칭",
    ], left=Inches(7.3), top=Inches(2.85),
       width=Inches(5.1), height=Inches(1.7), size=10, color=SUB)
    add_pill(s, "DRIVER  ·  드라이버", left=Inches(7.3), top=Inches(4.75), fill=DRIVER_CLR)
    add_bullets(s, [
        "고객 차량 운전 (대리·일일·공항·통근·이벤트·시외)",
        "차종: Sedan / SUV / MPV / Van",
        "면허 등급: SIM A · SIM A Umum",
        "driverServices: designated/daily/airport/hourly/commute/intercity/event",
    ], left=Inches(7.3), top=Inches(5.15),
       width=Inches(5.1), height=Inches(1.7), size=10, color=SUB)
    add_footer(s, 4, TOTAL)


def s05_screens_main():
    s = new_slide()
    add_header(s, "11  ·  CUSTOMER  ·  MAP", "화면 캡처 #1 — 메인 흐름", page_label="지도 2/12")
    captions = [
        ("기본 상태",
         "지도 + 바텀시트 Peek\nGET /partners (bbox)\nGET /reverse-geocode"),
        ("Peek — 필터 칩",
         "필터 칩 9종\nGET /activities (마스터)\nGET /service-types"),
        ("Peek — 스토리",
         "파트너 사진 가로 스크롤\n(검색 결과 상위 N건)"),
        ("Expanded — 리스트",
         "FlatList 정렬·페이지네이션\nGET /partners?cursor=…"),
    ]
    left0 = Inches(0.55)
    top   = Inches(1.55)
    gap   = Inches(0.30)
    for i, (title, body) in enumerate(captions):
        left = left0 + (PHONE_W + gap) * i
        add_phone_placeholder(s, left, top, label=title, sub=f"01_{i+1:02d}_*.png")
        add_phone_caption(s, left, top + PHONE_H + Inches(0.10),
                          title=title, body=body)
    add_text(s,
             "※ screenshots/map/ 폴더에 11_map_##_*.png 형식으로 저장",
             left=Inches(0.6), top=Inches(7.05),
             width=Inches(11), height=Inches(0.3), size=9, color=SUB)
    add_footer(s, 5, TOTAL)


def s06_screens_filters():
    s = new_slide()
    add_header(s, "11  ·  CUSTOMER  ·  MAP", "화면 캡처 #2 — 마커 / 필터 모달", page_label="지도 3/12")
    captions = [
        ("마커 선택",
         "파트너 미리보기 시트\nGET /partners/{id}\nPOST /events (marker_click)"),
        ("활동 필터",
         "다중 선택 (활동 → skill 매핑)\nGET /activities"),
        ("서비스 유형",
         "regular / onetime / live-in\n단일 선택"),
        ("조건 필터",
         "검증 토글 + 경력\n1y / 3y / 5y"),
    ]
    left0 = Inches(0.55)
    top   = Inches(1.55)
    gap   = Inches(0.30)
    for i, (title, body) in enumerate(captions):
        left = left0 + (PHONE_W + gap) * i
        add_phone_placeholder(s, left, top, label=title, sub=f"02_{i+1:02d}_*.png")
        add_phone_caption(s, left, top + PHONE_H + Inches(0.10),
                          title=title, body=body)
    add_footer(s, 6, TOTAL)


def s07_ui():
    s = new_slide()
    add_header(s, "11  ·  CUSTOMER  ·  MAP", "UI 구성 (영역별)", page_label="지도 4/12")
    cards = [
        ("상단 검색 바", [
            "파트너 이름 / 위치 키워드",
            "위치 키워드 → /geocode 호출",
            "디바운스 300ms · 최소 2자",
        ]),
        ("지도 (Google Maps)", [
            "커스텀 스타일 (밝은 회색 톤)",
            "POI / 대중교통 라벨 숨김",
            "현재 위치 표시 (권한 시)",
            "마커: 사진+이름+평점 · 헬퍼/드라이버 색상 구분",
            "줌 ± · 내 위치 버튼",
        ]),
        ("바텀시트 — Peek (240pt)", [
            "헤더: \"n명의 파트너\" 카운트",
            "필터 칩 9종 가로 스크롤",
            "스토리 리스트: 사진 가로 스크롤",
        ]),
        ("바텀시트 — Expanded (72%)", [
            "정렬 드롭다운",
            "파트너 카드 FlatList",
            "탭 → 헬퍼/드라이버 상세",
        ]),
        ("선택 미리보기 시트", [
            "사진·이름·평점·시급·위치",
            "스킬·경력·검증 배지",
            "\"상세 보기\" CTA",
        ]),
        ("필터 모달 4종", [
            "활동 카테고리 (다중)",
            "서비스 유형 (단일)",
            "헬퍼 조건 (검증+경력)",
            "정렬 (평점/후기/가격)",
        ]),
    ]
    cell_w = Inches(3.97)
    cell_h = Inches(2.62)
    gap_x  = Inches(0.16)
    gap_y  = Inches(0.18)
    left0  = Inches(0.6)
    top0   = Inches(1.55)
    for i, (title, items) in enumerate(cards):
        col = i % 3
        row = i // 3
        left = left0 + (cell_w + gap_x) * col
        top  = top0  + (cell_h + gap_y) * row
        add_rect(s, left, top, cell_w, cell_h, fill=BG_SOFT, line=LINE, radius=0.04)
        add_text(s, title, left=left + Inches(0.25), top=top + Inches(0.20),
                 width=cell_w - Inches(0.5), height=Inches(0.4),
                 size=12, bold=True, color=INK)
        add_bullets(s, items, left=left + Inches(0.25), top=top + Inches(0.6),
                    width=cell_w - Inches(0.5), height=cell_h - Inches(0.7),
                    size=10, color=SUB, line_spacing=1.3)
    add_footer(s, 7, TOTAL)


def s08_actions():
    s = new_slide()
    add_header(s, "11  ·  CUSTOMER  ·  MAP", "사용자 액션 → 백엔드 매핑", page_label="지도 5/12")

    rows = [
        ("액션", "트리거", "백엔드 호출"),
        ("진입",                 "탭 전환 / 라우트",          "GET /partners (bbox 기반)"),
        ("위치 권한 허용",       "최초 진입",                  "GET /reverse-geocode (현재 위치 → 동네)"),
        ("지도 이동·줌",         "팬·핀치",                    "GET /partners (debounce 400ms, 새 bbox)"),
        ("검색어 입력",          "TextInput 디바운스",         "GET /partners?keyword=…  +  /geocode"),
        ("마커 탭",              "Marker onPress",             "GET /partners/{id}  +  POST /events (marker_click)"),
        ("필터 적용",            "모달 \"적용\"",              "GET /partners (필터 파라미터 전체)"),
        ("정렬 변경",            "정렬 모달",                  "GET /partners?sortBy=…"),
        ("리스트 페이징",        "FlatList onEndReached",      "GET /partners?cursor=…"),
        ("상세 진입",            "\"상세 보기\" CTA",          "GET /partners/{id}/availability"),
        ("즐겨찾기 토글",        "하트 버튼 (상세시트)",       "POST/DELETE /favorites/partners/{id}"),
        ("신고",                 "메뉴 → 신고",                "POST /reports/partners/{id}"),
        ("차단",                 "메뉴 → 차단",                "POST /blocks/partners/{id}  → 검색 결과 필터링"),
    ]
    left = Inches(0.6); top = Inches(1.55)
    col_widths = [Inches(2.6), Inches(3.5), Inches(6.03)]
    row_h = Inches(0.36)

    for r_idx, row in enumerate(rows):
        is_header = r_idx == 0
        bg = INK if is_header else (BG_SOFT if r_idx % 2 == 1 else RGBColor(0xFF, 0xFF, 0xFF))
        cur_left = left
        for c_idx, cell in enumerate(row):
            add_rect(s, cur_left, top, col_widths[c_idx], row_h, fill=bg,
                     line=LINE if not is_header else None)
            add_text(s, cell,
                     left=cur_left + Inches(0.12), top=top + Inches(0.08),
                     width=col_widths[c_idx] - Inches(0.2), height=row_h,
                     size=10, bold=is_header,
                     color=RGBColor(0xFF, 0xFF, 0xFF) if is_header else INK,
                     font=MONO_FONT if (c_idx == 2 and not is_header) else KO_FONT)
            cur_left += col_widths[c_idx]
        top += row_h
    add_footer(s, 8, TOTAL)


def s09_data_partner():
    s = new_slide()
    add_header(s, "11  ·  CUSTOMER  ·  MAP", "데이터 모델 — Partner", page_label="지도 6/12")
    rows = [
        ("필드", "타입", "설명"),
        ("id",                "string",                            "파트너 고유 ID"),
        ("partnerType",       "'helper' | 'driver'",               "파트너 종류"),
        ("name / firstName",  "string",                            "표시명"),
        ("photoUrl / photoThumbUrl", "string",                     "원본 / 썸네일 분리 권장"),
        ("location",          "string",                            "표시용 동네명 (다국어 정책 필요)"),
        ("areaId",            "string (FK → Area)",                "활성 서비스 지역 ID"),
        ("lat, lng",          "number",                            "위치 좌표"),
        ("distanceMeters",    "number  (응답 시 계산)",            "사용자 위치 기준 — 거리 정렬용"),
        ("rating",            "number (0~5)",                      "평균 평점 (집계)"),
        ("reviewCount",       "number",                            "후기 수"),
        ("totalJobs",         "number",                            "누적 작업 건수"),
        ("pricePerHour",      "number (IDR)",                      "시급"),
        ("priceMin / priceMax", "number",                          "지역·시간대별 변동 가격 (정책에 따라)"),
        ("isAvailable",       "boolean",                           "즉시 가능 (스케줄/토글 정책)"),
        ("availabilityUntil", "ISO datetime",                      "현재 가용 만료 시점 (실시간 신호)"),
        ("isVerified",        "boolean",                           "신원·서류 검증"),
        ("verificationLevel", "'basic' | 'identity' | 'background'", "검증 단계 세분화 권장"),
        ("experienceYears",   "number",                            "경력 연수"),
        ("serviceFrequency",  "'regular'|'special'|'both'",        "정기/단기/둘다"),
        ("skills",            "string[]",                          "Beberes / Masak / Cuci / Setrika / …"),
        ("driverServices?",   "string[]",                          "designated/daily/airport/hourly/…"),
        ("vehicleTypes?",     "string[]",                          "Sedan/SUV/MPV/Van"),
        ("licenseClass?",     "string",                            "SIM A · SIM A Umum"),
        ("isFavorited",       "boolean",                           "현재 사용자 기준 즐겨찾기 여부"),
        ("isBlocked",         "boolean",                           "차단 시 검색에서 자동 제외 (서버측)"),
    ]
    left = Inches(0.6); top = Inches(1.50)
    col_widths = [Inches(2.85), Inches(3.5), Inches(5.78)]
    row_h = Inches(0.215)

    for r_idx, row in enumerate(rows):
        is_header = r_idx == 0
        bg = INK if is_header else (BG_SOFT if r_idx % 2 == 1 else RGBColor(0xFF, 0xFF, 0xFF))
        cur_left = left
        for c_idx, cell in enumerate(row):
            add_rect(s, cur_left, top, col_widths[c_idx], row_h, fill=bg,
                     line=LINE if not is_header else None)
            add_text(s, cell,
                     left=cur_left + Inches(0.10), top=top + Inches(0.025),
                     width=col_widths[c_idx] - Inches(0.18), height=row_h,
                     size=9, bold=is_header,
                     color=RGBColor(0xFF, 0xFF, 0xFF) if is_header else INK)
            cur_left += col_widths[c_idx]
        top += row_h
    add_footer(s, 9, TOTAL)


def s10_data_aux():
    s = new_slide()
    add_header(s, "11  ·  CUSTOMER  ·  MAP", "데이터 모델 — 보조 엔티티", page_label="지도 7/12")

    boxes = [
        ("Activity (활동 카테고리)", [
            "id : string  (cleaning, cooking, …)",
            "label : { id, en, ko, zh, ja }",
            "icon · color · bg",
            "items[] : ActivityItem",
        ]),
        ("ActivityItem (세부 항목)", [
            "id · label(다국어) · skill(매핑)",
            "예: { id: 'floor', skill: 'Beberes' }",
            "→ 검색 필터의 selectedActivities 와 매핑",
        ]),
        ("Skill (스킬)", [
            "code : string  (Beberes, Masak, …)",
            "label : { id, en, ko, … }",
            "category : 'cleaning'|'cooking'|…",
            "검색 필터·파트너 프로필 공용",
        ]),
        ("Area (활성 서비스 지역)", [
            "id · name (다국어) · centerLat · centerLng",
            "polygon (선택, 경계)",
            "isActive : boolean",
            "→ \"인근 지역\" 추천·검증",
        ]),
        ("Favorite", [
            "userId · partnerId · createdAt",
            "유니크 인덱스 (userId, partnerId)",
        ]),
        ("Report / Block", [
            "userId · partnerId · reason · createdAt",
            "Block은 검색에서 서버측 자동 제외",
            "Report는 운영자 큐에 누적",
        ]),
    ]
    cell_w = Inches(3.97)
    cell_h = Inches(2.62)
    gap_x  = Inches(0.16)
    gap_y  = Inches(0.18)
    left0  = Inches(0.6)
    top0   = Inches(1.55)
    for i, (title, items) in enumerate(boxes):
        col = i % 3
        row = i // 3
        left = left0 + (cell_w + gap_x) * col
        top  = top0  + (cell_h + gap_y) * row
        add_rect(s, left, top, cell_w, cell_h, fill=BG_SOFT, line=LINE, radius=0.04)
        add_text(s, title, left=left + Inches(0.25), top=top + Inches(0.20),
                 width=cell_w - Inches(0.5), height=Inches(0.4),
                 size=12, bold=True, color=INK)
        add_bullets(s, items, left=left + Inches(0.25), top=top + Inches(0.6),
                    width=cell_w - Inches(0.5), height=cell_h - Inches(0.7),
                    size=10, color=SUB, line_spacing=1.3)
    add_footer(s, 10, TOTAL)


def s11_api_search():
    s = new_slide()
    add_header(s, "11  ·  CUSTOMER  ·  MAP", "API #1 — 검색 / 상세 / 가용 시간", page_label="지도 8/12")

    add_api_card(
        s, left=Inches(0.6), top=Inches(1.55), width=Inches(7.5), height=Inches(5.45),
        method="GET", path="/api/v1/partners",
        title="파트너 검색 (메인)",
        body_lines=[
            "bbox=swLat,swLng,neLat,neLng    또는    lat / lng / radius",
            "partnerType: helper | driver | all",
            "serviceType: regular | onetime | live-in",
            "availableOnly · verifiedOnly  (boolean)",
            "minExperience: 0 | 1 | 3 | 5",
            "skills[]   ←  활동 → skill 매핑 결과",
            "keyword    ←  이름 또는 위치 키워드",
            "sortBy: rating | reviews | price_low | price_high | distance",
            "userLat / userLng  ←  거리 정렬·distanceMeters 계산용",
            "cursor / limit  (페이지네이션, cursor 권장)",
            "응답: { items[], nextCursor, total }",
            "  · items[i] = Partner (마커+카드 공용 필드)",
        ],
        method_color=BRAND,
        hint="성능: 마커 응답은 가벼워야 함. 매우 많을 경우 마커 전용 엔드포인트 분리 검토.",
    )

    add_api_card(
        s, left=Inches(8.30), top=Inches(1.55), width=Inches(4.43), height=Inches(2.62),
        method="GET", path="/partners/{id}",
        title="단일 파트너 상세",
        body_lines=[
            "미리보기 + 상세 화면 공용",
            "추가 필드: bio · workingHours · reviews",
            "isFavorited · verificationLevel",
        ],
        method_color=BRAND,
    )

    add_api_card(
        s, left=Inches(8.30), top=Inches(4.38), width=Inches(4.43), height=Inches(2.62),
        method="GET", path="/partners/{id}/availability",
        title="가용 시간 슬롯",
        body_lines=[
            "from / to  ←  조회 기간 (기본 7일)",
            "응답: 슬롯 배열 [{start, end, isAvailable}]",
            "예약 화면과 공용",
        ],
        method_color=BRAND,
    )
    add_footer(s, 11, TOTAL)


def s12_api_master_geo():
    s = new_slide()
    add_header(s, "11  ·  CUSTOMER  ·  MAP", "API #2 — 마스터 데이터 / 지오코딩", page_label="지도 9/12")

    cards = [
        dict(method="GET", path="/api/v1/activities",
             title="활동 카테고리 마스터",
             body_lines=[
                 "다국어 라벨 (id/ko/en/zh/ja)",
                 "category → item → skill 매핑",
                 "ETag/Cache-Control 권장",
                 "서버 관리 (앱 업데이트 없이 갱신)",
             ]),
        dict(method="GET", path="/api/v1/skills",
             title="스킬 마스터",
             body_lines=[
                 "code · label(다국어) · category",
                 "활동 매핑 결과의 단일 출처",
             ]),
        dict(method="GET", path="/api/v1/areas",
             title="활성 서비스 지역",
             body_lines=[
                 "id · name(다국어) · center · polygon?",
                 "검색·홈 화면 공용",
                 "신규 지역 출시 제어용",
             ]),
        dict(method="GET", path="/api/v1/geocode?q=…",
             title="지오코딩 (위치 키워드 → 좌표)",
             body_lines=[
                 "검색창 위치 입력 시",
                 "결과: [{name, lat, lng, areaId}]",
                 "외부 Provider 위임 가능 (Google/Mapbox)",
             ]),
        dict(method="GET", path="/reverse-geocode?lat&lng",
             title="역지오코딩 (좌표 → 동네명)",
             body_lines=[
                 "최초 진입 시 \"내 동네\" 표시",
                 "외부 Provider 캐시 권장",
             ]),
        dict(method="GET", path="/api/v1/pricing/policy",
             title="요금 정책 (선택)",
             body_lines=[
                 "지역·시간대별 요금 차등 시 노출",
                 "지도 카드의 priceMin/Max 산출 기준",
                 "정책 미적용 시 생략 가능",
             ]),
    ]
    cell_w = Inches(3.97)
    cell_h = Inches(2.62)
    gap_x  = Inches(0.16)
    gap_y  = Inches(0.18)
    left0  = Inches(0.6)
    top0   = Inches(1.55)
    for i, card in enumerate(cards):
        col = i % 3
        row = i // 3
        left = left0 + (cell_w + gap_x) * col
        top  = top0  + (cell_h + gap_y) * row
        add_api_card(s, left=left, top=top, width=cell_w, height=cell_h, **card)
    add_footer(s, 12, TOTAL)


def s13_api_user_actions():
    s = new_slide()
    add_header(s, "11  ·  CUSTOMER  ·  MAP", "API #3 — 사용자 행동 (즐겨찾기 / 신고 / 차단 / 추천)", page_label="지도 10/12")

    cards = [
        dict(method="POST", path="/favorites/partners/{id}",
             title="즐겨찾기 추가",
             body_lines=[
                 "유니크 (userId, partnerId)",
                 "Idempotent (이미 있으면 200)",
             ],
             method_color=INFO),
        dict(method="DEL", path="/favorites/partners/{id}",
             title="즐겨찾기 해제",
             body_lines=[
                 "없으면 204 (관대 처리)",
             ],
             method_color=WARN),
        dict(method="GET", path="/favorites/partners",
             title="내 즐겨찾기 목록",
             body_lines=[
                 "프로필 / 홈에서 공용",
                 "isFavorited 응답에 사전 반영",
             ],
             method_color=BRAND),
        dict(method="POST", path="/reports/partners/{id}",
             title="신고",
             body_lines=[
                 "reason · description · evidenceUrls[]",
                 "운영자 큐 적재",
             ],
             method_color=INFO),
        dict(method="POST", path="/blocks/partners/{id}",
             title="차단",
             body_lines=[
                 "이후 검색·추천에서 서버측 자동 제외",
                 "isBlocked 플래그 응답",
             ],
             method_color=INFO),
        dict(method="GET", path="/partners/recommended",
             title="추천 파트너 (개인화)",
             body_lines=[
                 "최근 본 / 즐겨찾기 / 검색 이력 기반",
                 "콜드 스타트: 평점·후기·인기 fallback",
                 "지도 진입 시 stories에 노출",
             ],
             method_color=BRAND,
             hint="추천 모델 위치(서버/외부)와 응답 SLA 합의 필요"),
    ]
    cell_w = Inches(3.97)
    cell_h = Inches(2.62)
    gap_x  = Inches(0.16)
    gap_y  = Inches(0.18)
    left0  = Inches(0.6)
    top0   = Inches(1.55)
    for i, card in enumerate(cards):
        col = i % 3
        row = i // 3
        left = left0 + (cell_w + gap_x) * col
        top  = top0  + (cell_h + gap_y) * row
        add_api_card(s, left=left, top=top, width=cell_w, height=cell_h, **card)
    add_footer(s, 13, TOTAL)


def s14_api_realtime_analytics():
    s = new_slide()
    add_header(s, "11  ·  CUSTOMER  ·  MAP", "API #4 — 실시간 / 검색 이력 / 분석", page_label="지도 11/12")

    cards = [
        dict(method="GET", path="/partners/recent",
             title="최근 본 파트너",
             body_lines=[
                 "스토리 리스트 / 홈 \"이어보기\"",
                 "서버측 자동 기록 (마커→상세 진입 시)",
             ],
             method_color=BRAND),
        dict(method="POST", path="/search/history",
             title="검색 이력 기록 (선택)",
             body_lines=[
                 "keyword · filters · resultCount",
                 "추천·자동완성 입력",
             ],
             method_color=INFO),
        dict(method="GET", path="/search/suggestions",
             title="자동완성·추천 키워드",
             body_lines=[
                 "최근 검색어 + 인기 키워드",
                 "검색창 포커스 시",
             ],
             method_color=BRAND),
        dict(method="WS", path="/ws/partners/availability",
             title="실시간 가용 상태",
             body_lines=[
                 "구독: { bbox, partnerType }",
                 "이벤트: { partnerId, isAvailable, availabilityUntil }",
                 "대안: 30초 폴링 — GET /partners/availability?ids=…",
             ],
             method_color=PURPLE,
             hint="WS 운영 부담 시 폴링부터. 트래픽·UX 따라 결정."),
        dict(method="POST", path="/events",
             title="분석 이벤트 로깅",
             body_lines=[
                 "type: map_view | marker_click | filter_apply | sort_change",
                 "props: { filters, resultCount, partnerId? }",
                 "→ 펀널 분석·검색 0건 추적",
             ],
             method_color=INFO),
        dict(method="POST", path="/devices/push-token",
             title="푸시 토큰 (간접)",
             body_lines=[
                 "지도 자체는 직접 사용 X",
                 "예약·검증 알림 위해 앱 시작 시 등록",
             ],
             method_color=INFO),
    ]
    cell_w = Inches(3.97)
    cell_h = Inches(2.62)
    gap_x  = Inches(0.16)
    gap_y  = Inches(0.18)
    left0  = Inches(0.6)
    top0   = Inches(1.55)
    for i, card in enumerate(cards):
        col = i % 3
        row = i // 3
        left = left0 + (cell_w + gap_x) * col
        top  = top0  + (cell_h + gap_y) * row
        add_api_card(s, left=left, top=top, width=cell_w, height=cell_h, **card)
    add_footer(s, 14, TOTAL)


def s15_policy_edge():
    s = new_slide()
    add_header(s, "11  ·  CUSTOMER  ·  MAP", "정책 합의 / 엣지 케이스 / 다음 단계", page_label="지도 12/12")

    add_text(s, "백엔드와 합의 필요", left=Inches(0.6), top=Inches(1.55),
             width=Inches(6), height=Inches(0.4), size=14, bold=True, color=WARN)
    add_bullets(s, [
        "isAvailable 정의 — 본인 토글 vs 스케줄 기반 vs 둘 다",
        "위치 갱신 주기 — 등록 거점 vs 실시간 GPS",
        "거리 정렬 도입 여부 / 거리 계산을 서버에서 할지",
        "verificationLevel 단계 정의 (basic/identity/background)",
        "헬퍼·드라이버 통합 검색 vs 분리 (스키마 갈라지면 분리)",
        "활동 → skill 매핑을 서버/클라 어디서?",
        "검색 빈도 / 캐싱 / Rate Limit 정책",
        "요금 정책 (단일 시급 vs 지역·시간대 차등)",
        "실시간 가용 상태: WS vs 폴링 (운영 부담 vs UX)",
        "차단된 파트너 서버측 자동 제외 보장",
        "location 다국어 표기 (사용자 lang 기반 변환?)",
        "추천(/partners/recommended) 모델 위치·SLA",
    ], left=Inches(0.6), top=Inches(2.0),
       width=Inches(6.0), height=Inches(4.7), size=10, color=INK, line_spacing=1.45)

    add_text(s, "엣지 케이스", left=Inches(7.0), top=Inches(1.55),
             width=Inches(6), height=Inches(0.4), size=14, bold=True, color=BRAND)
    add_bullets(s, [
        "위치 권한 거부 → 기본 데모 영역",
        "결과 0건 → 빈 상태 + 필터 초기화 (미구현)",
        "네트워크 에러 → 재시도 토스트 (미구현)",
        "비검증 파트너 → 배지 미노출",
        "즉시 불가 → \"예약 불가\" 표기, CTA 비활성",
        "마커 100+ → 클러스터링 검토",
        "서비스 미진입 지역 → \"인근 활성 지역 안내\" 컴포넌트",
        "키워드 매칭 0건 + 위치도 0건 → 추천으로 fallback",
    ], left=Inches(7.0), top=Inches(2.0),
       width=Inches(5.7), height=Inches(3.5), size=10, color=INK, line_spacing=1.45)

    add_text(s, "다음 단계", left=Inches(7.0), top=Inches(5.55),
             width=Inches(6), height=Inches(0.4), size=14, bold=True, color=PURPLE)
    add_bullets(s, [
        "지도 페이지 폰 스크린샷 8장 캡처 → 슬라이드 5·6에 삽입",
        "백엔드와 본 슬라이드 12 \"합의 필요\" 항목 워크숍",
        "다음 화면 명세 진행 (홈 / 헬퍼 검색 / 회원가입 중 택)",
    ], left=Inches(7.0), top=Inches(5.95),
       width=Inches(5.7), height=Inches(1.2), size=10, color=SUB, line_spacing=1.4)

    add_footer(s, 15, TOTAL)


# ── 빌드 ─────────────────────────────────────
TOTAL = 15
s01_cover()
s02_index()
s03_common()
s04_map_overview()
s05_screens_main()
s06_screens_filters()
s07_ui()
s08_actions()
s09_data_partner()
s10_data_aux()
s11_api_search()
s12_api_master_geo()
s13_api_user_actions()
s14_api_realtime_analytics()
s15_policy_edge()

out = "/Users/bigband/Linka-app/Linka-app/docs/backend-spec/Linka-backend-spec-v2.pptx"
prs.save(out)
print(f"Saved: {out}  ({TOTAL} slides)")
