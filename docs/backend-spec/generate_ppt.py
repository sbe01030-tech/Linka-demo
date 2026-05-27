"""
Linka 백엔드 기획안 PPT 생성기.
실행: /tmp/linka-ppt/bin/python docs/backend-spec/generate_ppt.py
출력: docs/backend-spec/Linka-backend-spec.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Constants ────────────────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

KO_FONT = "Apple SD Gothic Neo"
EN_FONT = "SF Pro Text"

BRAND       = RGBColor(0x00, 0xC8, 0x53)   # green
INK         = RGBColor(0x10, 0x14, 0x18)
SUB         = RGBColor(0x55, 0x60, 0x6C)
LINE        = RGBColor(0xE3, 0xE6, 0xEA)
BG_SOFT     = RGBColor(0xF5, 0xF7, 0xF8)
HELPER_CLR  = RGBColor(0xF5, 0x9E, 0x0B)
DRIVER_CLR  = RGBColor(0x63, 0x66, 0xF1)
WARN        = RGBColor(0xEF, 0x44, 0x44)


def set_font(run, *, size=14, bold=False, color=INK, font=KO_FONT):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_text(slide, text, *, left, top, width, height,
             size=14, bold=False, color=INK, align=PP_ALIGN.LEFT, font=KO_FONT):
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_font(run, size=size, bold=bold, color=color, font=font)
    return tx


def add_bullets(slide, items, *, left, top, width, height,
                size=12, color=INK, line_spacing=1.35):
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = "•  " + item
        set_font(run, size=size, color=color)
    return tx


def add_rect(slide, left, top, width, height, *, fill=BG_SOFT, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.shadow.inherit = False
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.5)
    return shape


def add_pill(slide, text, *, left, top, fill, color=RGBColor(0xFF, 0xFF, 0xFF), size=10):
    width = Inches(1.4)
    height = Inches(0.32)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.shadow.inherit = False
    shape.adjustments[0] = 0.5
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    set_font(run, size=size, bold=True, color=color)
    return shape


def add_header(slide, eyebrow, title, *, page_label=""):
    # eyebrow
    add_text(slide, eyebrow, left=Inches(0.6), top=Inches(0.45),
             width=Inches(8), height=Inches(0.3), size=11, bold=True, color=BRAND)
    # title
    add_text(slide, title, left=Inches(0.6), top=Inches(0.78),
             width=Inches(11), height=Inches(0.7), size=26, bold=True, color=INK)
    # divider
    add_rect(slide, Inches(0.6), Inches(1.55), Inches(12.13), Emu(9525), fill=LINE)
    # page label (top right)
    if page_label:
        add_text(slide, page_label, left=Inches(11.2), top=Inches(0.5),
                 width=Inches(1.5), height=Inches(0.3), size=10, color=SUB,
                 align=PP_ALIGN.RIGHT)


def add_footer(slide, page_no, total):
    add_text(slide, "Linka  ·  Backend Spec  ·  v0.1",
             left=Inches(0.6), top=Inches(7.05),
             width=Inches(6), height=Inches(0.3), size=9, color=SUB)
    add_text(slide, f"{page_no} / {total}",
             left=Inches(11.2), top=Inches(7.05),
             width=Inches(1.5), height=Inches(0.3), size=9, color=SUB,
             align=PP_ALIGN.RIGHT)


def add_screenshot_placeholder(slide, left, top, width, height, label):
    box = add_rect(slide, left, top, width, height, fill=BG_SOFT, line=LINE)
    add_text(slide, "📷 스크린샷 삽입", left=left, top=top + Inches(0.25),
             width=width, height=Inches(0.4), size=12, color=SUB,
             align=PP_ALIGN.CENTER)
    add_text(slide, label, left=left, top=top + Inches(0.7),
             width=width, height=Inches(0.4), size=10, color=SUB,
             align=PP_ALIGN.CENTER)


# ── Build presentation ──────────────────────────────────────
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]


def new_slide():
    s = prs.slides.add_slide(blank_layout)
    bg = add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=RGBColor(0xFF, 0xFF, 0xFF))
    return s


# 임시 — 정확한 페이지 번호는 마지막에 후처리하지 않고 직접 카운트
SLIDES = []


def s_cover():
    s = new_slide()
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=RGBColor(0xFA, 0xFB, 0xFC))
    # 좌측 강조 띠
    add_rect(s, 0, 0, Inches(0.35), SLIDE_H, fill=BRAND)
    add_text(s, "BACKEND SPEC", left=Inches(0.9), top=Inches(1.6),
             width=Inches(8), height=Inches(0.4), size=12, bold=True, color=BRAND)
    add_text(s, "Linka 백엔드 기획안", left=Inches(0.9), top=Inches(2.0),
             width=Inches(11), height=Inches(1.0), size=44, bold=True, color=INK)
    add_text(s,
             "React Native(Expo) 기반 인도네시아 타겟 가사·운전·심부름 중개 앱",
             left=Inches(0.9), top=Inches(3.1),
             width=Inches(11), height=Inches(0.5), size=16, color=SUB)
    add_text(s, "v0.1  ·  2026-04-27", left=Inches(0.9), top=Inches(6.6),
             width=Inches(6), height=Inches(0.3), size=11, color=SUB)
    SLIDES.append(("cover", s))


def s_index():
    s = new_slide()
    add_header(s, "01  ·  INDEX", "전체 화면 목차")
    cols = [
        ("인증 (Auth)", [
            "01  스플래시", "02  환영 화면 (역할 선택)",
            "03  로그인", "04  회원가입", "05  약관",
        ]),
        ("고객 (Customer)", [
            "10  홈", "11  지도 (파트너 탐색) ← 우선",
            "12  헬퍼 검색", "13  헬퍼 상세",
            "14  예약", "15  주문 내역",
            "16  리뷰 작성", "17  프로필", "18  알림",
        ]),
        ("헬퍼 / 드라이버 / 심부름", [
            "20  헬퍼 홈", "21  헬퍼 주문", "22  헬퍼 프로필",
            "30  드라이버 게시판", "31  드라이버 상세",
            "40  심부름 게시판", "41  심부름 등록", "42  심부름 상세",
        ]),
        ("채팅 / 공통", [
            "50  채팅 목록", "51  채팅 상세",
            "60  프로필 수정", "61  도움말 / FAQ",
            "",
            "※ 커뮤니티 탭은 본 기획안에서 제외",
        ]),
    ]
    col_w = Inches(2.95)
    gap   = Inches(0.15)
    left0 = Inches(0.6)
    top0  = Inches(1.85)
    for i, (title, items) in enumerate(cols):
        left = left0 + (col_w + gap) * i
        add_text(s, title, left=left, top=top0,
                 width=col_w, height=Inches(0.4), size=13, bold=True, color=BRAND)
        add_bullets(s, items, left=left, top=top0 + Inches(0.45),
                    width=col_w, height=Inches(4.5), size=11, color=INK)
    add_footer(s, 2, TOTAL)
    SLIDES.append(("index", s))


def s_common():
    s = new_slide()
    add_header(s, "02  ·  COMMON", "공통 사항")
    # 4분할 카드
    cards = [
        ("사용자 역할", [
            "customer  · 서비스 요청자",
            "helper      · 가사도우미 (ART)",
            "driver       · 운전자 (대리/일일)",
            "errand      · 심부름꾼 (예정)",
            "→ 멀티 역할 허용 여부 정책 결정 필요",
        ]),
        ("다국어", [
            "지원: id (기본) · en · ko · zh · ja",
            "사용자 입력은 원문 저장",
            "카테고리·시스템 메시지만 번역 키 보유",
        ]),
        ("통화 / 지역", [
            "통화: IDR (Rp) · 정수 단위 저장",
            "기본 지역: 자카르타 (Jakarta)",
            "지도 중심: Kemang / Kebayoran Baru",
            "좌표: WGS84 (lat/lng)",
        ]),
        ("인증 / 미디어", [
            "이메일 또는 전화번호 (정책 필요)",
            "JWT 기반 세션",
            "약관 동의 이력 보관",
            "사진·첨부: Pre-signed URL 권장",
        ]),
    ]
    box_w = Inches(5.95)
    box_h = Inches(2.45)
    positions = [
        (Inches(0.6),  Inches(1.85)),
        (Inches(6.78), Inches(1.85)),
        (Inches(0.6),  Inches(4.45)),
        (Inches(6.78), Inches(4.45)),
    ]
    for (title, items), (left, top) in zip(cards, positions):
        add_rect(s, left, top, box_w, box_h, fill=BG_SOFT, line=LINE)
        add_text(s, title, left=left + Inches(0.3), top=top + Inches(0.25),
                 width=box_w - Inches(0.6), height=Inches(0.4),
                 size=14, bold=True, color=INK)
        add_bullets(s, items, left=left + Inches(0.3), top=top + Inches(0.7),
                    width=box_w - Inches(0.6), height=box_h - Inches(0.85),
                    size=11, color=SUB)
    add_footer(s, 3, TOTAL)
    SLIDES.append(("common", s))


# ── 지도 화면 ──────────────────────────────
def s_map_overview():
    s = new_slide()
    add_header(s, "11  ·  CUSTOMER  ·  MAP", "지도 화면 — 개요", page_label="지도 1/6")
    # 좌: 목적·진입
    add_text(s, "목적", left=Inches(0.6), top=Inches(1.85),
             width=Inches(6), height=Inches(0.4), size=14, bold=True, color=BRAND)
    add_text(s,
             "지도 위에서 주변 헬퍼·드라이버를 시각적으로 탐색하고, 필터·정렬을 통해 원하는 파트너를 찾아 상세 진입한다.",
             left=Inches(0.6), top=Inches(2.25),
             width=Inches(6), height=Inches(1.2), size=13, color=INK)

    add_text(s, "진입 경로", left=Inches(0.6), top=Inches(3.55),
             width=Inches(6), height=Inches(0.4), size=14, bold=True, color=BRAND)
    add_bullets(s, [
        "하단 탭바 → \"지도(Map)\"",
        "홈 화면 → 서비스 카드 / \"지도에서 보기\"",
        "라우트 파라미터:  expanded?: boolean",
        "라우트 파라미터:  serviceType?: regular | onetime | live-in",
    ], left=Inches(0.6), top=Inches(3.95),
       width=Inches(6), height=Inches(2.5), size=12, color=INK)

    # 우: 파트너 종류 비주얼
    add_rect(s, Inches(7.0), Inches(1.85), Inches(5.7), Inches(4.6),
             fill=BG_SOFT, line=LINE)
    add_text(s, "두 종류의 파트너", left=Inches(7.3), top=Inches(2.05),
             width=Inches(5.1), height=Inches(0.4), size=14, bold=True, color=INK)

    add_pill(s, "HELPER  ·  헬퍼", left=Inches(7.3), top=Inches(2.6), fill=HELPER_CLR)
    add_bullets(s, [
        "가사도우미 (ART)",
        "스킬: Beberes / Masak / Cuci / Setrika / Deep Cleaning",
        "정기 / 단기 / 상주",
        "검증·경력·평점 기반 매칭",
    ], left=Inches(7.3), top=Inches(2.95),
       width=Inches(5.1), height=Inches(1.5), size=11, color=SUB)

    add_pill(s, "DRIVER  ·  드라이버", left=Inches(7.3), top=Inches(4.65), fill=DRIVER_CLR)
    add_bullets(s, [
        "고객 차량 운전 (대리·일일·공항·통근 등)",
        "차종: Sedan / SUV / MPV / Van",
        "면허 등급: SIM A · SIM A Umum",
    ], left=Inches(7.3), top=Inches(5.0),
       width=Inches(5.1), height=Inches(1.3), size=11, color=SUB)

    add_footer(s, 4, TOTAL)
    SLIDES.append(("map_overview", s))


def s_map_screens():
    s = new_slide()
    add_header(s, "11  ·  CUSTOMER  ·  MAP", "지도 화면 — 화면 캡처", page_label="지도 2/6")
    labels = [
        "01_default — 기본 상태 (지도 + peek)",
        "02_expanded — 바텀시트 펼침",
        "03_marker_selected — 마커 선택",
        "04_filter_activity — 활동 필터",
        "05_filter_service — 서비스 유형",
        "06_filter_condition — 조건 필터",
    ]
    # 3열 × 2행
    cell_w = Inches(3.97)
    cell_h = Inches(2.45)
    gap_x  = Inches(0.16)
    gap_y  = Inches(0.18)
    left0  = Inches(0.6)
    top0   = Inches(1.85)
    for i, label in enumerate(labels):
        col = i % 3
        row = i // 3
        left = left0 + (cell_w + gap_x) * col
        top  = top0  + (cell_h + gap_y) * row
        add_screenshot_placeholder(s, left, top, cell_w, cell_h, label)
    add_text(s, "※ screenshots/map/ 폴더에 같은 파일명으로 저장 후 끼워넣기",
             left=Inches(0.6), top=Inches(7.0),
             width=Inches(11), height=Inches(0.3), size=10, color=SUB)
    add_footer(s, 5, TOTAL)
    SLIDES.append(("map_screens", s))


def s_map_ui():
    s = new_slide()
    add_header(s, "11  ·  CUSTOMER  ·  MAP", "지도 화면 — UI 구성", page_label="지도 3/6")
    cards = [
        ("상단 검색 바", [
            "파트너 이름 / 위치 키워드",
            "클리어(X) 버튼",
        ]),
        ("지도 (Google Maps)", [
            "커스텀 스타일 (밝은 회색, POI/대중교통 숨김)",
            "현재 위치 표시 (권한 시)",
            "마커: 사진 + 이름 + 평점, 헬퍼/드라이버 색상 구분",
            "줌 ± · 내 위치 버튼",
        ]),
        ("바텀시트 — Peek (240pt)", [
            "헤더: \"n명의 파트너\" 카운트",
            "필터 칩 가로 스크롤 (9종)",
            "스토리 리스트: 사진 가로 스크롤",
        ]),
        ("바텀시트 — Expanded (72%)", [
            "정렬 드롭다운",
            "파트너 카드 FlatList",
            "카드 탭 → 헬퍼/드라이버 상세 이동",
        ]),
        ("선택 미리보기 시트", [
            "사진·이름·평점·시급·위치·스킬·경력·검증 배지",
            "\"상세 보기\" CTA",
        ]),
        ("필터 모달 4종", [
            "활동 카테고리 (다중 선택)",
            "서비스 유형 (단일)",
            "헬퍼 조건 (검증 + 경력)",
            "정렬 (평점/후기/가격)",
        ]),
    ]
    cell_w = Inches(3.97)
    cell_h = Inches(2.55)
    gap_x  = Inches(0.16)
    gap_y  = Inches(0.18)
    left0  = Inches(0.6)
    top0   = Inches(1.85)
    for i, (title, items) in enumerate(cards):
        col = i % 3
        row = i // 3
        left = left0 + (cell_w + gap_x) * col
        top  = top0  + (cell_h + gap_y) * row
        add_rect(s, left, top, cell_w, cell_h, fill=BG_SOFT, line=LINE)
        add_text(s, title, left=left + Inches(0.25), top=top + Inches(0.2),
                 width=cell_w - Inches(0.5), height=Inches(0.4),
                 size=12, bold=True, color=INK)
        add_bullets(s, items, left=left + Inches(0.25), top=top + Inches(0.6),
                    width=cell_w - Inches(0.5), height=cell_h - Inches(0.7),
                    size=10, color=SUB, line_spacing=1.3)
    add_footer(s, 6, TOTAL)
    SLIDES.append(("map_ui", s))


def s_map_data():
    s = new_slide()
    add_header(s, "11  ·  CUSTOMER  ·  MAP", "지도 화면 — Partner 데이터 모델", page_label="지도 4/6")

    # 테이블
    rows = [
        ("필드", "타입", "설명"),
        ("id",                "string",                            "파트너 고유 ID"),
        ("name / firstName",  "string",                            "표시명"),
        ("partnerType",       "'helper' | 'driver'",               "파트너 종류"),
        ("lat, lng",          "number",                            "위치 좌표"),
        ("photo",             "string (URL)",                      "프로필 사진"),
        ("location",          "string",                            "표시용 동네명 (예: Kebayoran Baru)"),
        ("rating",            "number (0~5)",                      "평균 평점 (집계)"),
        ("totalJobs",         "number",                            "누적 작업 건수"),
        ("pricePerHour",      "number (IDR)",                      "시급"),
        ("isAvailable",       "boolean",                           "즉시 가능 여부"),
        ("isVerified",        "boolean",                           "신원·서류 검증"),
        ("experienceYears",   "number",                            "경력 연수"),
        ("serviceFrequency",  "'regular'|'special'|'both'",        "정기/단기/둘다"),
        ("skills",            "string[]",                          "Beberes / Masak / Cuci / Setrika / …"),
        ("driverServices?",   "string[]  (driver only)",           "designated/daily/airport/hourly/…"),
        ("licenseClass?",     "string  (driver only)",             "SIM A · SIM A Umum"),
    ]
    left = Inches(0.6); top = Inches(1.85)
    col_widths = [Inches(2.6), Inches(3.3), Inches(6.23)]
    row_h = Inches(0.32)

    for r_idx, row in enumerate(rows):
        is_header = r_idx == 0
        bg = INK if is_header else (BG_SOFT if r_idx % 2 == 1 else RGBColor(0xFF, 0xFF, 0xFF))
        cur_left = left
        for c_idx, cell in enumerate(row):
            add_rect(s, cur_left, top, col_widths[c_idx], row_h, fill=bg,
                     line=LINE if not is_header else None)
            add_text(s, cell,
                     left=cur_left + Inches(0.12), top=top + Inches(0.05),
                     width=col_widths[c_idx] - Inches(0.2), height=row_h,
                     size=10, bold=is_header,
                     color=RGBColor(0xFF, 0xFF, 0xFF) if is_header else INK)
            cur_left += col_widths[c_idx]
        top += row_h

    add_footer(s, 7, TOTAL)
    SLIDES.append(("map_data", s))


def s_map_api():
    s = new_slide()
    add_header(s, "11  ·  CUSTOMER  ·  MAP", "지도 화면 — API 요구사항", page_label="지도 5/6")

    # 카드 1: 메인 검색
    add_rect(s, Inches(0.6), Inches(1.85), Inches(7.7), Inches(5.0),
             fill=BG_SOFT, line=LINE)
    add_pill(s, "GET",
             left=Inches(0.85), top=Inches(2.05),
             fill=BRAND, size=10)
    add_text(s, "/api/v1/partners",
             left=Inches(2.4), top=Inches(2.05),
             width=Inches(5.5), height=Inches(0.4), size=14, bold=True, color=INK,
             font=EN_FONT)
    add_text(s, "파트너 검색 (메인)", left=Inches(0.85), top=Inches(2.45),
             width=Inches(7.0), height=Inches(0.3), size=11, color=SUB)
    add_bullets(s, [
        "bbox=swLat,swLng,neLat,neLng    또는   lat/lng/radius",
        "partnerType: helper | driver | all",
        "serviceType: regular | onetime | live-in",
        "availableOnly · verifiedOnly  (boolean)",
        "minExperience: 0 | 1 | 3 | 5",
        "skills[] · keyword",
        "sortBy: rating | reviews | price_low | price_high",
        "cursor / limit  (페이지네이션)",
    ], left=Inches(0.85), top=Inches(2.85),
       width=Inches(7.2), height=Inches(3.9), size=10, color=INK, line_spacing=1.35)

    # 카드 2: 상세 + 카테고리
    add_rect(s, Inches(8.5), Inches(1.85), Inches(4.23), Inches(2.4),
             fill=BG_SOFT, line=LINE)
    add_pill(s, "GET", left=Inches(8.75), top=Inches(2.05), fill=BRAND, size=10)
    add_text(s, "/api/v1/partners/{id}",
             left=Inches(10.3), top=Inches(2.05),
             width=Inches(2.5), height=Inches(0.4), size=12, bold=True, color=INK,
             font=EN_FONT)
    add_text(s, "단일 파트너 상세", left=Inches(8.75), top=Inches(2.45),
             width=Inches(4.0), height=Inches(0.3), size=11, color=SUB)
    add_bullets(s, [
        "미리보기 + 상세 화면 공용",
        "자기소개·후기·가능 시간대 포함",
    ], left=Inches(8.75), top=Inches(2.85),
       width=Inches(3.7), height=Inches(1.3), size=10, color=INK)

    add_rect(s, Inches(8.5), Inches(4.45), Inches(4.23), Inches(2.4),
             fill=BG_SOFT, line=LINE)
    add_pill(s, "GET", left=Inches(8.75), top=Inches(4.65), fill=BRAND, size=10)
    add_text(s, "/api/v1/activities",
             left=Inches(10.3), top=Inches(4.65),
             width=Inches(2.5), height=Inches(0.4), size=12, bold=True, color=INK,
             font=EN_FONT)
    add_text(s, "활동 카테고리 마스터", left=Inches(8.75), top=Inches(5.05),
             width=Inches(4.0), height=Inches(0.3), size=11, color=SUB)
    add_bullets(s, [
        "다국어 라벨 (id/ko/en/zh/ja)",
        "category → item → skill 매핑",
        "서버 관리 권장 (앱 업데이트 없이 갱신)",
    ], left=Inches(8.75), top=Inches(5.45),
       width=Inches(3.7), height=Inches(1.3), size=10, color=INK)

    add_footer(s, 8, TOTAL)
    SLIDES.append(("map_api", s))


def s_map_policy():
    s = new_slide()
    add_header(s, "11  ·  CUSTOMER  ·  MAP", "지도 화면 — 합의 필요 / 미해결", page_label="지도 6/6")

    add_text(s, "백엔드와 합의 필요", left=Inches(0.6), top=Inches(1.85),
             width=Inches(6), height=Inches(0.4), size=14, bold=True, color=WARN)
    add_bullets(s, [
        "isAvailable 정의 — 본인 토글 vs 스케줄 기반",
        "위치 갱신 주기 — 등록 거점 vs 실시간 GPS",
        "거리 정렬 도입 여부",
        "isVerified 기준 (서류·절차)",
        "헬퍼·드라이버 통합 검색 vs 분리",
        "활동 → skill 매핑을 서버/클라 어디서?",
        "검색 빈도 / 캐싱 정책",
        "location 다국어 표기 정책",
    ], left=Inches(0.6), top=Inches(2.3),
       width=Inches(6.0), height=Inches(4.6), size=11, color=INK, line_spacing=1.45)

    add_text(s, "엣지 케이스", left=Inches(7.0), top=Inches(1.85),
             width=Inches(6), height=Inches(0.4), size=14, bold=True, color=BRAND)
    add_bullets(s, [
        "위치 권한 거부 → 기본 데모 영역",
        "결과 0건 → 빈 상태 + 필터 초기화 (미구현)",
        "네트워크 에러 → 재시도 토스트 (미구현)",
        "비검증 파트너 → 배지 미노출",
        "즉시 불가 → \"예약 불가\" 표기, CTA 비활성",
        "마커 100+ → 클러스터링 검토",
    ], left=Inches(7.0), top=Inches(2.3),
       width=Inches(5.7), height=Inches(3.5), size=11, color=INK, line_spacing=1.45)

    add_text(s, "다음 단계", left=Inches(7.0), top=Inches(5.95),
             width=Inches(6), height=Inches(0.4), size=14, bold=True, color=BRAND)
    add_bullets(s, [
        "스크린샷 7장 캡처 → 슬라이드 5에 삽입",
        "다음 화면(홈 / 회원가입 등) 명세 진행",
    ], left=Inches(7.0), top=Inches(6.4),
       width=Inches(5.7), height=Inches(1.0), size=11, color=SUB)

    add_footer(s, 9, TOTAL)
    SLIDES.append(("map_policy", s))


# ── 슬라이드 빌드 ───────────────────────────
TOTAL = 9  # cover + index + common + 6 map slides
s_cover()
s_index()
s_common()
s_map_overview()
s_map_screens()
s_map_ui()
s_map_data()
s_map_api()
s_map_policy()

out = "/Users/bigband/Linka-app/Linka-app/docs/backend-spec/Linka-backend-spec.pptx"
prs.save(out)
print(f"Saved: {out}  ({len(SLIDES)} slides)")
