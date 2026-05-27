"""
Linka 철학 슬라이드 — 한 페이지 (16:9).
슬로건: "결국 사람. 옮기는 온기."

실행: /tmp/linka-ppt/bin/python docs/philosophy/generate_philosophy.py
출력: docs/philosophy/Linka-philosophy.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── 상수 ────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

BRAND  = RGBColor(0x00, 0xC8, 0x53)
INK    = RGBColor(0x10, 0x14, 0x18)
SUB    = RGBColor(0x55, 0x60, 0x6C)
SOFT   = RGBColor(0xA8, 0xB0, 0xB7)
LINE   = RGBColor(0xE3, 0xE6, 0xEA)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BG     = RGBColor(0xFC, 0xFD, 0xFD)

KO_FONT     = "Apple SD Gothic Neo"
EN_FONT     = "Inter"
BRAND_FONT  = "Nunito"

HERE     = os.path.dirname(os.path.abspath(__file__))
LOGO     = os.path.normpath(os.path.join(HERE, "..", "..", "assets", "logo", "linka_circle_2048.png"))
OUT_PATH = os.path.join(HERE, "Linka-philosophy.pptx")


# ── 유틸 ────────────────────────────────────
def set_font(run, *, size, bold=False, color=INK, font=KO_FONT, italic=False):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def add_text(slide, text, *, left, top, width, height,
             size, bold=False, color=INK, align=PP_ALIGN.LEFT,
             font=KO_FONT, line_spacing=None, italic=False):
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing is not None:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    set_font(run, size=size, bold=bold, color=color, font=font, italic=italic)
    return tx


def add_rect(slide, left, top, width, height, *, fill=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.shadow.inherit = False
    if fill is not None:
        shape.fill.solid(); shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    shape.line.fill.background()
    return shape


def add_oval(slide, cx, cy, r, *, fill=None, line_color=None, line_width=None):
    """중심·반지름 기준 동그라미."""
    left = cx - r
    top = cy - r
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, r * 2, r * 2)
    shape.shadow.inherit = False
    if fill is not None:
        shape.fill.solid(); shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line_color is not None:
        shape.line.color.rgb = line_color
        if line_width is not None:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


# ── 슬라이드 빌드 ─────────────────────────────
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
slide = prs.slides.add_slide(prs.slide_layouts[6])

# 배경
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)

# ── 상단 헤더 ────────────────────────────────
# 좌측: Linka 워드마크
add_text(slide, "Linka",
         left=Inches(0.7), top=Inches(0.55),
         width=Inches(2), height=Inches(0.45),
         size=20, bold=True, color=BRAND, font=BRAND_FONT)

# 우측: 작은 라벨 (OUR PHILOSOPHY · 철학)
add_text(slide, "OUR PHILOSOPHY  ·  철학",
         left=Inches(8), top=Inches(0.7),
         width=Inches(4.7), height=Inches(0.3),
         size=10, bold=True, color=SUB, font=EN_FONT,
         align=PP_ALIGN.RIGHT)

# 헤더 구분선 (얇고 옅음)
add_rect(slide, Inches(0.7), Inches(1.35),
         Inches(12), Emu(9525), fill=LINE)

# ── 본문 좌측 ───────────────────────────────
# 슬로건 좌측 — 작은 그린 수직 바 (인용 라인 느낌)
add_rect(slide, Inches(0.7), Inches(2.55),
         Inches(0.06), Inches(2.5), fill=BRAND)

# 슬로건 1줄 — "결국 사람."  (사람 = 그린)
slogan1 = slide.shapes.add_textbox(
    Inches(1.0), Inches(2.4), Inches(8.0), Inches(1.4))
tf = slogan1.text_frame
tf.margin_left = 0; tf.margin_right = 0
tf.margin_top = 0;  tf.margin_bottom = 0
p = tf.paragraphs[0]
p.line_spacing = 1.0
for txt, color in [("결국 ", INK), ("사람", BRAND), (".", INK)]:
    r = p.add_run()
    r.text = txt
    set_font(r, size=68, bold=True, color=color, font=KO_FONT)

# 슬로건 2줄 — "옮기는 온기."  (온기 = 그린)
slogan2 = slide.shapes.add_textbox(
    Inches(1.0), Inches(3.55), Inches(8.0), Inches(1.4))
tf2 = slogan2.text_frame
tf2.margin_left = 0; tf2.margin_right = 0
tf2.margin_top = 0;  tf2.margin_bottom = 0
p2 = tf2.paragraphs[0]
p2.line_spacing = 1.0
for txt, color in [("옮기는 ", INK), ("온기", BRAND), (".", INK)]:
    r = p2.add_run()
    r.text = txt
    set_font(r, size=68, bold=True, color=color, font=KO_FONT)

# 영문 서브타이틀
add_text(slide, "In the end, people. Carrying warmth.",
         left=Inches(1.0), top=Inches(4.9),
         width=Inches(8), height=Inches(0.4),
         size=14, color=SUB, font=EN_FONT, italic=True)

# 짧은 그린 강조 라인 (구분)
add_rect(slide, Inches(1.0), Inches(5.55),
         Inches(0.5), Emu(38100), fill=BRAND)

# 본문 단락 (2줄)
body = slide.shapes.add_textbox(
    Inches(1.0), Inches(5.8), Inches(7.5), Inches(1.4))
tf3 = body.text_frame
tf3.margin_left = 0; tf3.margin_right = 0
tf3.margin_top = 0;  tf3.margin_bottom = 0
tf3.word_wrap = True

p3 = tf3.paragraphs[0]
p3.line_spacing = 1.55
r = p3.add_run()
r.text = "기술은 거리를 좁힙니다. 하지만 사람을 잇는 건, 결국 사람의 손길입니다."
set_font(r, size=12, color=SUB, font=KO_FONT)

p4 = tf3.add_paragraph()
p4.line_spacing = 1.55
r = p4.add_run()
r.text = "Linka는 도움이 필요한 가정과 도울 수 있는 사람을 잇는, 가장 따뜻한 다리입니다."
set_font(r, size=12, color=SUB, font=KO_FONT)

# ── 우측 시각 요소 ──────────────────────────
# 옅은 그린 데코 원 (아주 큰 outline 원 — 배경 액센트)
deco_cx = Inches(10.8)
deco_cy = Inches(3.9)
add_oval(slide, deco_cx, deco_cy + Inches(0.3), Inches(2.55),
         fill=None, line_color=BRAND, line_width=Pt(1.0))
# 더 작고 옅은 원
add_oval(slide, deco_cx - Inches(1.8), deco_cy + Inches(1.4), Inches(0.35),
         fill=BRAND, line_color=None)
add_oval(slide, deco_cx + Inches(1.7), deco_cy - Inches(1.0), Inches(0.18),
         fill=BRAND, line_color=None)

# 메인 — Linka 로고 뱃지 (그린 원 + 마스코트 + Linka)
if os.path.exists(LOGO):
    slide.shapes.add_picture(
        LOGO,
        Inches(9.0), Inches(2.1),
        width=Inches(3.6), height=Inches(3.6),
    )

# ── 하단 푸터 ────────────────────────────────
# 옅은 하단 구분선
add_rect(slide, Inches(0.7), Inches(7.0),
         Inches(12), Emu(9525), fill=LINE)

add_text(slide, "© 2026 Linka  ·  Small links. Warmer world.",
         left=Inches(0.7), top=Inches(7.15),
         width=Inches(8), height=Inches(0.3),
         size=9, color=SOFT, font=EN_FONT)

add_text(slide, "01 / 01",
         left=Inches(11.0), top=Inches(7.15),
         width=Inches(1.6), height=Inches(0.3),
         size=9, color=SOFT, font=EN_FONT,
         align=PP_ALIGN.RIGHT)

# ── 저장 ────────────────────────────────────
prs.save(OUT_PATH)
print(f"Saved: {OUT_PATH}")
